import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors

# --- Configuration ---
OUTPUT_DIR = r"c:\Users\Hemanth\Downloads\Hemanth\ML Project\documentation"
os.makedirs(OUTPUT_DIR, exist_ok=True)

PPT_PATH = os.path.join(OUTPUT_DIR, "Project_Presentation.pptx")
PDF_PATH = os.path.join(OUTPUT_DIR, "Project_Logic_Explanation.pdf")

# --- Content Data ---
CONTENT = {
    "title": "AI Podcast Generator",
    "subtitle": "Intelligent NLP & Neural Text-to-Speech Pipeline",
    "overview": [
        "End-to-end automated system transforming topics into high-quality audio podcasts.",
        "Leverages Generative AI for both content creation and voice synthesis.",
        "Focuses on natural, conversational dialogue between a Host and a Guest."
    ],
    "architecture": [
        "1. Script Generation: User Prompt -> Gemini 2.5 Flash -> JSON Script",
        "2. Audio Generation: JSON Script -> Edge TTS -> Audio Chunks",
        "3. Audio Processing: Audio Chunks -> FFmpeg -> Final MP3"
    ],
    "ml_details": {
        "NLP": [
            "Model: Google Gemini 2.5 Flash",
            "Strategy: Zero-Shot scriptwriting with persona constraints",
            "Output: Structured JSON using Pydantic validation"
        ],
        "TTS": [
            "Model: Microsoft Azure Neural Voices (Edge TTS)",
            "Voices: en-US-ChristopherNeural & en-US-JennyNeural",
            "Logic: Asynchronous synthesis of dialogue segments"
        ]
    },
    "tech_stack": [
        "Python (Streamlit, Pydantic, Asyncio)",
        "Gemini API (GenAI)",
        "Edge-TTS (Azure Neural Engine)",
        "FFmpeg (Audio Processing Engine)"
    ]
}

# --- PDF Generation ---
def generate_pdf():
    doc = SimpleDocTemplate(PDF_PATH, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []

    # Title
    styles.add(ParagraphStyle(name='CenterTitle', parent=styles['Title'], alignment=1))
    story.append(Paragraph(CONTENT["title"], styles['CenterTitle']))
    story.append(Paragraph(CONTENT["subtitle"], styles['Italic']))
    story.append(Spacer(1, 12))

    # Overview
    story.append(Paragraph("1. Project Overview", styles['Heading1']))
    for line in CONTENT["overview"]:
        story.append(Paragraph(f"• {line}", styles['Normal']))
    story.append(Spacer(1, 12))

    # Architecture
    story.append(Paragraph("2. System Architecture", styles['Heading1']))
    for line in CONTENT["architecture"]:
        story.append(Paragraph(line, styles['Normal']))
    story.append(Spacer(1, 12))

    # ML Details
    story.append(Paragraph("3. Machine Learning Components", styles['Heading1']))
    
    story.append(Paragraph("A. Scriptwriting (NLP)", styles['Heading2']))
    for line in CONTENT["ml_details"]["NLP"]:
        story.append(Paragraph(f"  - {line}", styles['Normal']))
        
    story.append(Paragraph("B. Audio Synthesis (TTS)", styles['Heading2']))
    for line in CONTENT["ml_details"]["TTS"]:
        story.append(Paragraph(f"  - {line}", styles['Normal']))
    story.append(Spacer(1, 12))

    # Tech Stack
    story.append(Paragraph("4. Tech Stack", styles['Heading1']))
    story.append(Paragraph(", ".join(CONTENT["tech_stack"]), styles['Normal']))

    doc.build(story)
    print(f"PDF generated: {PDF_PATH}")

# --- PPT Generation ---
def generate_ppt():
    prs = Presentation()
    
    # helper for dark theme
    def apply_dark_theme(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(15, 15, 35) # Dark Blue/Navy

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    apply_dark_theme(slide)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = CONTENT["title"]
    subtitle.text = CONTENT["subtitle"]
    
    # Style Title
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title.text_frame.paragraphs[0].font.size = Pt(44)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)

    # Slide 2: Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_dark_theme(slide)
    slide.shapes.title.text = "Project Overview"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 255, 255) # Cyan
    
    body = slide.shapes.placeholders[1].text_frame
    for line in CONTENT["overview"]:
        p = body.add_paragraph()
        p.text = line
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.level = 0

    # Slide 3: Architecture (The Pipeline)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_dark_theme(slide)
    slide.shapes.title.text = "System Architecture"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 255) # Magenta
    
    body = slide.shapes.placeholders[1].text_frame
    for line in CONTENT["architecture"]:
        p = body.add_paragraph()
        p.text = line
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.level = 0

    # Slide 4: ML - NLP
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_dark_theme(slide)
    slide.shapes.title.text = "Logic: Script Generation (NLP)"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 0) # Yellow
    
    body = slide.shapes.placeholders[1].text_frame
    for line in CONTENT["ml_details"]["NLP"]:
        p = body.add_paragraph()
        p.text = line
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.level = 0

    # Slide 5: ML - TTS
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_dark_theme(slide)
    slide.shapes.title.text = "Logic: Audio Synthesis (TTS)"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 255, 0) # Green
    
    body = slide.shapes.placeholders[1].text_frame
    for line in CONTENT["ml_details"]["TTS"]:
        p = body.add_paragraph()
        p.text = line
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.level = 0

    # Slide 6: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_dark_theme(slide)
    slide.shapes.title.text = "Tech Stack Summary"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    body = slide.shapes.placeholders[1].text_frame
    p = body.add_paragraph()
    p.text = ", ".join(CONTENT["tech_stack"])
    p.font.color.rgb = RGBColor(200, 200, 200)

    prs.save(PPT_PATH)
    print(f"PPT generated: {PPT_PATH}")

if __name__ == "__main__":
    generate_pdf()
    generate_ppt()
